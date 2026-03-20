"""Advanced AI PPT generation: charts, flowcharts, icons, auto-layout."""
import io
import json
import re
import uuid
import logging
from datetime import datetime
from typing import Optional, List, Dict, Any

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select

from ..db.database import get_db
from ..db.models import Project
from ..core.gemini import generate_text

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/slides", tags=["slides"])

# In-memory store (upgrade to DB later)
_store: Dict[str, dict] = {}

# ── Request / Response Models ─────────────────────────────

class GenerateRequest(BaseModel):
    topic: str
    num_slides: int = 10
    template: str = "professional"
    extra_context: str = ""
    file_ids: List[str] = []


# ── Templates ─────────────────────────────────────────

THEMES = {
    "professional": {
        "bg":         (15,  15,  30),
        "header_bg":  (10,  10,  20),
        "chapter_bg": (30,  30,  60),
        "accent":     (99,  102, 241),
        "title_fg":   (255, 255, 255),
        "body_fg":    (200, 200, 220),
        "dark":       True,
    },
    "modern": {
        "bg":         (248, 249, 255),
        "header_bg":  (99,  102, 241),
        "chapter_bg": (224, 225, 255),
        "accent":     (99,  102, 241),
        "title_fg":   (255, 255, 255),
        "body_fg":    (40,  40,  70),
        "dark":       False,
    },
    "minimal": {
        "bg":         (255, 255, 255),
        "header_bg":  (248, 248, 252),
        "chapter_bg": (243, 244, 255),
        "accent":     (99,  102, 241),
        "title_fg":   (20,  20,  40),
        "body_fg":    (60,  60,  80),
        "dark":       False,
    },
}

CHART_PALETTE = [
    "#6366f1", "#8b5cf6", "#ec4899", "#f59e0b",
    "#10b981", "#3b82f6", "#ef4444", "#14b8a6",
]

# ── AI Prompt ───────────────────────────────────────────

PROMPT_TMPL = """你是頂級的專業簡報設計師與數據分析師。請根據主題生成一份完整、高品質的繁體中文簡報。

主題：{topic}
投影片數量：約 {num_slides} 張
{extra}

要求：
1. 自動分章節（每章節用 chapter 類型投影片開頭）
2. 根據主題選擇最適合的版型組合
3. 數據要合理真實，圖表數據要有意義
4. content 的每個要點必須選用最貼切的 emoji 圖示
5. 必須包含至少 1 個圖表（bar_chart 或 pie_chart）
6. 若主題有流程/步驟，必須包含 flowchart
7. 第一張必須是 title，最後一張必須是 summary

嚴格以 JSON 格式回傳（不加任何說明文字）：
{{
  "title": "簡報主標題",
  "subtitle": "副標題",
  "slides": [
    {{"id":"1","type":"title","title":"主標題","subtitle":"副標題","notes":""}},
    {{"id":"2","type":"chapter","title":"第一章","subtitle":"章節說明","icon":"🎯","notes":""}},
    {{"id":"3","type":"content","title":"投影片標題","content":[{{"icon":"✅","text":"要點一詳細說明"}},{{"icon":"📊","text":"要點二詳細說明"}},{{"icon":"🚀","text":"要點三詳細說明"}},{{"icon":"💡","text":"要點四詳細說明"}}],"notes":"演講備注"}},
    {{"id":"4","type":"two_column","title":"對比分析","left_title":"優勢","left_content":["A點","B點","C點"],"right_title":"挑戰","right_content":["X點","Y點","Z點"],"notes":""}},
    {{"id":"5","type":"bar_chart","title":"數據比較","chart":{{"labels":["類別1","類別2","類別3","類別4","類別5"],"values":[45,72,58,89,63],"unit":"單位","color":"indigo"}},"notes":"圖表說明"}},
    {{"id":"6","type":"pie_chart","title":"比例分析","chart":{{"labels":["項目A","項目B","項目C","項目D"],"values":[35,28,22,15],"unit":"%"}},"notes":""}},
    {{"id":"7","type":"line_chart","title":"趨勢分析","chart":{{"labels":["1月","2月","3月","4月","5月","6月"],"values":[30,45,42,67,58,80],"unit":"單位"}},"notes":""}},
    {{"id":"8","type":"flowchart","title":"流程圖","flow":{{"nodes":[{{"id":"n1","label":"開始","type":"start"}},{{"id":"n2","label":"步驟一","type":"process"}},{{"id":"n3","label":"判斷？","type":"decision"}},{{"id":"n4","label":"步驟二","type":"process"}},{{"id":"n5","label":"完成","type":"end"}}],"connections":[{{"from_":"n1","to":"n2","label":""}},{{"from_":"n2","to":"n3","label":""}},{{"from_":"n3","to":"n4","label":"是"}},{{"from_":"n3","to":"n2","label":"否"}},{{"from_":"n4","to":"n5","label":""}}]}},"notes":"流程說明"}},
    {{"id":"9","type":"table","title":"數據總覽","table":{{"headers":["項目","數値","說明"],"rows":[["項目1","100","說明1"],["項目2","200","說明2"],["項目3","150","說明3"]]}},"notes":""}},
    {{"id":"10","type":"quote","title":"","quote":"這裡是最重要的核心洞見或關鍵引述，讓觀眾銘記於心。","author":"資料來源","notes":""}},
    {{"id":"11","type":"summary","title":"總結","content":["核心要點一","核心要點二","核心要點三"],"cta":"立即行動","notes":""}}
  ]
}}

根據主題靈活調整，不必每種類型都用，選最適合的組合。"""

# ── Routes ────────────────────────────────────────────

@router.post("/{project_id}/generate")
async def generate_slides(
    project_id: str,
    req: GenerateRequest,
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(Project).where(Project.id == project_id))
    project = result.scalar_one_or_none()
    if not project:
        raise HTTPException(404, "Project not found")

    extra = f"額外背景資訊：{req.extra_context}" if req.extra_context else ""
    prompt = PROMPT_TMPL.format(
        topic=req.topic,
        num_slides=req.num_slides,
        extra=extra,
    )

    try:
        raw = await generate_text(prompt)
        m = re.search(r'\{[\s\S]*\}', raw)
        data = json.loads(m.group() if m else raw)
    except Exception as e:
        logger.error(f"AI generation error: {e}\nRaw: {raw[:500] if 'raw' in dir() else 'N/A'}")
        data = {
            "title": req.topic,
            "subtitle": "",
            "slides": [
                {"id": "1", "type": "title", "title": req.topic, "subtitle": "AI 生成失敗，請重試", "notes": ""},
            ],
        }

    pres_id = str(uuid.uuid4())
    pres = {
        "id": pres_id,
        "project_id": project_id,
        "topic": req.topic,
        "template": req.template,
        "title": data.get("title", req.topic),
        "subtitle": data.get("subtitle", ""),
        "slides": data.get("slides", []),
        "created_at": datetime.utcnow().isoformat(),
    }
    _store[pres_id] = pres
    return pres


@router.get("/{project_id}")
async def list_slides(project_id: str):
    items = [
        {
            "id": p["id"], "title": p["title"], "topic": p["topic"],
            "template": p["template"],
            "slide_count": len(p["slides"]),
            "created_at": p["created_at"],
        }
        for p in _store.values() if p["project_id"] == project_id
    ]
    return sorted(items, key=lambda x: x["created_at"], reverse=True)


@router.get("/{project_id}/{pres_id}")
async def get_slides(project_id: str, pres_id: str):
    p = _store.get(pres_id)
    if not p or p["project_id"] != project_id:
        raise HTTPException(404)
    return p


@router.patch("/{project_id}/{pres_id}")
async def update_slides(project_id: str, pres_id: str, data: dict):
    p = _store.get(pres_id)
    if not p or p["project_id"] != project_id:
        raise HTTPException(404)
    p.update(data)
    return p


@router.delete("/{project_id}/{pres_id}")
async def delete_slides(project_id: str, pres_id: str):
    _store.pop(pres_id, None)
    return {"ok": True}


@router.get("/{project_id}/{pres_id}/download")
async def download_pptx(project_id: str, pres_id: str):
    p = _store.get(pres_id)
    if not p or p["project_id"] != project_id:
        raise HTTPException(404)
    pptx_bytes = _build_pptx(p)
    safe_title = re.sub(r'[^\w\u4e00-\u9fff\-_]', '_', p["title"])[:50]
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{safe_title}.pptx"'},
    )


# ── PPTX Builder ──────────────────────────────────────────

def _build_pptx(pres: dict) -> bytes:
    from pptx import Presentation as PPT
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    theme = THEMES.get(pres.get("template", "professional"), THEMES["professional"])

    p = PPT()
    p.slide_width  = Inches(13.33)
    p.slide_height = Inches(7.5)

    dispatch = {
        "title":      _slide_title,
        "chapter":    _slide_chapter,
        "content":    _slide_content,
        "two_column": _slide_two_col,
        "bar_chart":  _slide_bar,
        "line_chart": _slide_line,
        "pie_chart":  _slide_pie,
        "flowchart":  _slide_flow,
        "table":      _slide_table,
        "quote":      _slide_quote,
        "summary":    _slide_summary,
    }

    for s in pres.get("slides", []):
        fn = dispatch.get(s.get("type", "content"), _slide_content)
        fn(p, s, theme)

    buf = io.BytesIO()
    p.save(buf)
    return buf.getvalue()


# ── Helper ────────────────────────────────────────────

def _blank(p):
    return p.slides.add_slide(p.slide_layouts[6])

def _bg(slide, rgb):
    from pptx.dml.color import RGBColor
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)

def _rect(slide, l, t, w, h, fill=None, line=None, lw=1):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(*fill)
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = RGBColor(*line); sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def _txt(slide, l, t, w, h, text, sz, bold=False, italic=False,
         color=(255,255,255), align="left"):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    par = tf.paragraphs[0]
    run = par.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = RGBColor(*color)
    par.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    return tb

def _header(slide, title, theme):
    _bg(slide, theme["bg"])
    _rect(slide, 0, 0, 13.33, 1.15, fill=theme["header_bg"])
    _rect(slide, 0, 1.15, 0.08, 6.35, fill=theme["accent"])
    tc = theme["title_fg"] if theme.get("dark") or theme["header_bg"] == theme["bg"] else (255,255,255)
    if not theme.get("dark") and theme["header_bg"] != theme["bg"]:
        tc = (255,255,255)
    _txt(slide, 0.35, 0.18, 12.5, 0.85, title, 26, bold=True, color=tc)

def _chart_img(fig, bg_color) -> bytes:
    import matplotlib; matplotlib.use("Agg"); import matplotlib.pyplot as plt
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", facecolor=bg_color, dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf

# ── Slide Types ───────────────────────────────────────────

def _slide_title(p, s, theme):
    slide = _blank(p); _bg(slide, theme["header_bg"])
    _rect(slide, 0, 3.0, 13.33, 0.06, fill=theme["accent"])
    _txt(slide, 1, 1.6, 11.33, 1.5, s.get("title",""), 42, bold=True,
         color=theme["title_fg"], align="center")
    if s.get("subtitle"):
        _txt(slide, 1, 3.3, 11.33, 0.8, s["subtitle"], 22,
             color=(180,180,210), align="center")
    _rect(slide, 5.3, 6.7, 0.35, 0.35, fill=theme["accent"])
    _rect(slide, 5.85, 6.7, 0.35, 0.35, fill=tuple(min(255,c+50) for c in theme["accent"]))
    _rect(slide, 6.4, 6.7, 0.35, 0.35, fill=tuple(max(0,c-50) for c in theme["accent"]))

def _slide_chapter(p, s, theme):
    slide = _blank(p); _bg(slide, theme["chapter_bg"])
    _rect(slide, 0, 0, 0.18, 7.5, fill=theme["accent"])
    icon = s.get("icon","◆")
    _txt(slide, 0.5, 2.3, 1.8, 1.2, icon, 48, align="center")
    _txt(slide, 2.5, 2.5, 10.0, 1.0, s.get("title",""), 34, bold=True,
         color=theme["title_fg"] if theme.get("dark") else (20,20,50))
    if s.get("subtitle"):
        _txt(slide, 2.5, 3.7, 10.0, 0.7, s["subtitle"], 20,
             color=tuple(c+50 if c<200 else c for c in theme["accent"]))

def _slide_content(p, s, theme):
    slide = _blank(p)
    _header(slide, s.get("title",""), theme)
    items = s.get("content", [])
    if not items: return
    n = len(items)
    y0, step = 1.4, min(1.0, 5.6 / max(n, 1))
    for i, item in enumerate(items):
        if isinstance(item, dict):
            icon, text = item.get("icon","▸"), item.get("text","")
        else:
            icon, text = "▸", str(item)
        y = y0 + i * step
        _txt(slide, 0.35, y, 0.65, 0.75, icon, 20, color=theme["accent"])
        _txt(slide, 1.1,  y, 11.8, 0.75, text, 18, color=theme["body_fg"])

def _slide_two_col(p, s, theme):
    slide = _blank(p)
    _header(slide, s.get("title",""), theme)
    _rect(slide, 6.55, 1.3, 0.06, 5.9, fill=theme["accent"])
    accent2 = tuple(min(255,c+80) for c in theme["accent"])
    if s.get("left_title"):
        _txt(slide, 0.35, 1.35, 5.8, 0.55, s["left_title"], 16, bold=True, color=accent2)
    for i, item in enumerate(s.get("left_content",[])[:5]):
        _txt(slide, 0.35, 2.05+i*0.78, 5.8, 0.7, f"▸  {item}", 15, color=theme["body_fg"])
    if s.get("right_title"):
        _txt(slide, 6.85, 1.35, 6.0, 0.55, s["right_title"], 16, bold=True, color=accent2)
    for i, item in enumerate(s.get("right_content",[])[:5]):
        _txt(slide, 6.85, 2.05+i*0.78, 6.0, 0.7, f"▸  {item}", 15, color=theme["body_fg"])

def _make_bar_chart(labels, values, unit, theme, line=False):
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np
    dark = theme.get("dark", True)
    bg = tuple(c/255 for c in theme["bg"])
    fg = "white" if dark else "#1a1a2e"
    accent_hex = "#{:02x}{:02x}{:02x}".format(*theme["accent"])
    colors = [CHART_PALETTE[i % len(CHART_PALETTE)] for i in range(len(labels))]
    fig, ax = plt.subplots(figsize=(10, 4.8))
    fig.patch.set_facecolor(bg)
    ax.set_facecolor(tuple(min(1.0, c+0.04) for c in bg))
    if line:
        ax.plot(labels, values, color=accent_hex, linewidth=2.5, marker='o',
                markersize=7, markerfacecolor=accent_hex)
        ax.fill_between(range(len(labels)), values, alpha=0.15, color=accent_hex)
    else:
        bars = ax.bar(labels, values, color=colors, width=0.6, zorder=3)
        for bar, val in zip(bars, values):
            ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+max(values)*0.01,
                    f"{val}{unit}", ha="center", va="bottom", color=fg,
                    fontsize=9, fontweight="bold")
    ax.set_ylabel(unit, color=fg, fontsize=10)
    ax.tick_params(colors=fg, labelsize=10)
    ax.spines[["top","right","left"]].set_visible(False)
    ax.spines["bottom"].set_color(fg)
    ax.grid(axis="y", alpha=0.15, color=fg, zorder=0)
    plt.tight_layout()
    return fig, bg

def _slide_bar(p, s, theme):
    slide = _blank(p)
    _header(slide, s.get("title",""), theme)
    c = s.get("chart",{}); labels=c.get("labels",[]); values=c.get("values",[]); unit=c.get("unit","")
    if not (labels and values): return
    fig, bg = _make_bar_chart(labels, values, unit, theme, line=False)
    img = _chart_img(fig, bg)
    from pptx.util import Inches
    slide.shapes.add_picture(img, Inches(0.9), Inches(1.3), Inches(11.4), Inches(5.8))

def _slide_line(p, s, theme):
    slide = _blank(p)
    _header(slide, s.get("title",""), theme)
    c = s.get("chart",{}); labels=c.get("labels",[]); values=c.get("values",[]); unit=c.get("unit","")
    if not (labels and values): return
    fig, bg = _make_bar_chart(labels, values, unit, theme, line=True)
    img = _chart_img(fig, bg)
    from pptx.util import Inches
    slide.shapes.add_picture(img, Inches(0.9), Inches(1.3), Inches(11.4), Inches(5.8))

def _slide_pie(p, s, theme):
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    slide = _blank(p)
    _header(slide, s.get("title",""), theme)
    c = s.get("chart",{}); labels=c.get("labels",[]); values=c.get("values",[])
    if not (labels and values): return
    dark = theme.get("dark", True)
    bg = tuple(c2/255 for c2 in theme["bg"])
    fg = "white" if dark else "#1a1a2e"
    colors = [CHART_PALETTE[i % len(CHART_PALETTE)] for i in range(len(labels))]
    fig, ax = plt.subplots(figsize=(9, 5.2))
    fig.patch.set_facecolor(bg); ax.set_facecolor(bg)
    wedges, texts, autotexts = ax.pie(
        values, labels=labels, autopct="%1.1f%%", colors=colors,
        startangle=90, wedgeprops={"edgecolor": bg, "linewidth":2},
        textprops={"color": fg, "fontsize": 12}
    )
    for at in autotexts:
        at.set_color("white"); at.set_fontweight("bold"); at.set_fontsize(11)
    plt.tight_layout()
    img = _chart_img(fig, bg)
    from pptx.util import Inches
    slide.shapes.add_picture(img, Inches(1.8), Inches(1.2), Inches(9.6), Inches(6.0))

def _slide_flow(p, s, theme):
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    from matplotlib.patches import FancyBboxPatch
    slide = _blank(p)
    _header(slide, s.get("title",""), theme)
    flow = s.get("flow", {})
    nodes = flow.get("nodes", [])
    conns = flow.get("connections", [])
    if not nodes: return

    dark = theme.get("dark", True)
    bg = tuple(c/255 for c in theme["bg"])
    fg = "white" if dark else "#1a1a2e"
    accent = "#{:02x}{:02x}{:02x}".format(*theme["accent"])

    n = len(nodes)
    fig, ax = plt.subplots(figsize=(11, 4.5))
    fig.patch.set_facecolor(bg); ax.set_facecolor(bg)
    ax.set_xlim(-0.8, n - 0.2); ax.set_ylim(-1.2, 1.8); ax.axis("off")

    pos = {}
    for i, node in enumerate(nodes):
        x = i; y = 0; pos[node["id"]] = (x, y)
        ntype = node.get("type","process")
        label = node["label"]
        if ntype in ("start","end"):
            ellipse = mpatches.Ellipse((x,y), 0.9, 0.55, facecolor=accent, edgecolor="none", zorder=3)
            ax.add_patch(ellipse)
        elif ntype == "decision":
            diamond = plt.Polygon([[x,y+0.3],[x+0.5,y],[x,y-0.3],[x-0.5,y]],
                                  facecolor=accent, edgecolor="none", zorder=3)
            ax.add_patch(diamond)
        else:
            rect = FancyBboxPatch((x-0.45, y-0.28), 0.9, 0.56,
                                  boxstyle="round,pad=0.04",
                                  facecolor=accent, edgecolor="none", zorder=3)
            ax.add_patch(rect)
        ax.text(x, y, label, ha="center", va="center", color="white",
                fontsize=8, fontweight="bold", zorder=4)

    for conn in conns:
        x1, y1 = pos.get(conn.get("from_",""), (0,0))
        x2, y2 = pos.get(conn.get("to",""), (0,0))
        ax.annotate("", xy=(x2-0.46, y2), xytext=(x1+0.46, y1),
                    arrowprops=dict(arrowstyle="->", color=fg, lw=1.5))
        if conn.get("label"):
            ax.text((x1+x2)/2, (y1+y2)/2+0.22, conn["label"],
                    ha="center", color=fg, fontsize=7)

    plt.tight_layout()
    img = _chart_img(fig, bg)
    from pptx.util import Inches
    slide.shapes.add_picture(img, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.8))

def _slide_table(p, s, theme):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    slide = _blank(p)
    _header(slide, s.get("title",""), theme)
    tbl = s.get("table", {}); headers=tbl.get("headers",[]); rows=tbl.get("rows",[])
    if not headers: return
    ncols = len(headers); nrows = len(rows)+1
    row_h = min(0.55, 5.6 / nrows)
    shape = slide.shapes.add_table(nrows, ncols,
        Inches(0.7), Inches(1.4), Inches(11.9), Inches(row_h*nrows))
    t = shape.table
    for j, h in enumerate(headers):
        cell = t.cell(0,j); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(*theme["accent"])
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
        cell.text_frame.paragraphs[0].font.size = Pt(13)
    for i, row in enumerate(rows):
        for j, val in enumerate(row[:ncols]):
            cell = t.cell(i+1,j); cell.text = str(val)
            if i%2==1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(
                    min(255,theme["bg"][0]+20),
                    min(255,theme["bg"][1]+20),
                    min(255,theme["bg"][2]+30))
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme["body_fg"])
            cell.text_frame.paragraphs[0].font.size = Pt(12)

def _slide_quote(p, s, theme):
    slide = _blank(p); _bg(slide, theme["header_bg"])
    _txt(slide, 0.3, 0.2, 2.5, 2.5, "\u201C", 90,
         color=tuple(min(255,c+100) for c in theme["accent"]))
    _txt(slide, 1.2, 1.4, 10.8, 3.5, s.get("quote",""), 28, italic=True,
         color=theme["title_fg"], align="center")
    _rect(slide, 5.5, 5.7, 2.33, 0.06, fill=theme["accent"])
    if s.get("author"):
        _txt(slide, 0.5, 5.9, 12.3, 0.5, f"— {s['author']}", 15,
             color=(160,160,200), align="center")

def _slide_summary(p, s, theme):
    slide = _blank(p)
    _header(slide, s.get("title","總結"), theme)
    items = s.get("content", [])
    accent2 = tuple(min(255,c+80) for c in theme["accent"])
    for i, item in enumerate(items[:5]):
        _rect(slide, 0.35, 1.45+i*0.9, 0.06, 0.65, fill=theme["accent"])
        _txt(slide, 0.6, 1.45+i*0.9, 12.0, 0.75, item, 20, bold=True, color=accent2)
    if s.get("cta"):
        _rect(slide, 2.5, 6.3, 8.33, 0.85, fill=theme["accent"])
        _txt(slide, 2.5, 6.3, 8.33, 0.85, f"→  {s['cta']}", 20, bold=True,
             color=(255,255,255), align="center")
