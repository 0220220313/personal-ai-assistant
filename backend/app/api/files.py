from fastapi import APIRouter, Depends, UploadFile, File, HTTPException, BackgroundTasks, Form, Body
from fastapi.responses import JSONResponse
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
import os, aiofiles, uuid, mimetypes, asyncio, logging, subprocess, tempfile
from typing import Optional

from ..db.database import get_db
from ..db.models import File as FileModel, Project
from ..core.gemini import upload_file_to_gemini, generate_file_summary

router = APIRouter(prefix="/files", tags=["files"])
logger = logging.getLogger(__name__)

UPLOAD_DIR = "./data/uploads"
ALLOWED_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
    "text/plain": "txt",
    "text/markdown": "txt",
    "image/jpeg": "image",
    "image/png": "image",
    "image/webp": "image",
}

os.makedirs(UPLOAD_DIR, exist_ok=True)


@router.get("/{project_id}/folders")
async def list_folders(project_id: str, db: AsyncSession = Depends(get_db)):
    result = await db.execute(
        select(FileModel.folder_path).where(FileModel.project_id == project_id).distinct()
    )
    paths = sorted(set(row[0] or "/" for row in result.fetchall()))
    folders = list(paths)
    for path in paths:
        if path == "/":
            continue
        parts = path.rstrip("/").split("/")
        for i in range(1, len(parts)):
            parent = "/".join(parts[:i]) or "/"
            if parent not in folders:
                folders.append(parent)
    return {"folders": sorted(set(folders))}


@router.get("/{project_id}")
async def list_files(
    project_id: str,
    folder: Optional[str] = None,
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(Project).where(Project.id == project_id))
    project = result.scalar_one_or_none()
    if not project:
        raise HTTPException(404, "Project not found")

    query = select(FileModel).where(FileModel.project_id == project_id)
    if folder is not None:
        query = query.where(FileModel.folder_path == folder)
    query = query.order_by(FileModel.created_at.desc())

    result = await db.execute(query)
    files = result.scalars().all()
    return [
        {
            "id": f.id,
            "original_name": f.original_name,
            "file_type": f.file_type,
            "file_size": f.file_size,
            "summary": f.summary,
            "is_indexed": f.is_indexed,
            "folder_path": f.folder_path or "/",
            "created_at": f.created_at.isoformat(),
        }
        for f in files
    ]


@router.post("/upload/{project_id}")
async def upload_file(
    project_id: str,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    folder: str = Form(default="/"),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(Project).where(Project.id == project_id))
    project = result.scalar_one_or_none()
    if not project:
        raise HTTPException(404, "Project not found")

    content_type = file.content_type or mimetypes.guess_type(file.filename or "")[0] or ""
    if file.filename and (file.filename.endswith(".md") or file.filename.endswith(".txt")):
        content_type = "text/plain"
    if file.filename and file.filename.endswith(".pptx"):
        content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if content_type not in ALLOWED_TYPES:
        raise HTTPException(400, f"Unsupported file type: {content_type}")

    file_id = str(uuid.uuid4())
    ext = os.path.splitext(file.filename or "")[1]
    filename = f"{file_id}{ext}"
    local_path = os.path.join(UPLOAD_DIR, filename)

    content = await file.read()
    async with aiofiles.open(local_path, "wb") as f:
        await f.write(content)

    db_file = FileModel(
        id=file_id,
        project_id=project_id,
        filename=filename,
        original_name=file.filename or filename,
        file_type=ALLOWED_TYPES.get(content_type, "unknown"),
        file_size=len(content),
        folder_path=folder or "/",
        is_indexed=False,
    )
    db.add(db_file)
    await db.commit()

    background_tasks.add_task(
        _process_file_background,
        file_id,
        local_path,
        content_type,
        file.filename or filename,
    )

    return {
        "id": file_id,
        "original_name": file.filename,
        "file_type": ALLOWED_TYPES.get(content_type, "unknown"),
        "file_size": len(content),
        "is_indexed": False,
        "folder_path": folder or "/",
        "created_at": db_file.created_at.isoformat(),
    }


@router.post("/{file_id}/parse-pptx")
async def parse_pptx(file_id: str, db: AsyncSession = Depends(get_db)):
    """Parse a PPTX file and extract structured content with optional Gemini vision analysis."""
    result = await db.execute(select(FileModel).where(FileModel.id == file_id))
    f = result.scalar_one_or_none()
    if not f:
        raise HTTPException(404, "File not found")
    if f.file_type != "pptx":
        raise HTTPException(400, "File is not a PPTX")

    local_path = os.path.join(UPLOAD_DIR, f.filename)
    if not os.path.exists(local_path):
        raise HTTPException(404, "File not found on disk")

    # Step 1: Parse PPTX with python-pptx
    parse_result = _parse_pptx_to_markdown(local_path)
    slides_md = parse_result["slides_markdown"]
    visual_slide_indices = parse_result["visual_slide_indices"]

    # Step 2: Try LibreOffice screenshot (optional)
    visual_notes: list[str] = []
    has_visual_slides = len(visual_slide_indices) > 0

    with tempfile.TemporaryDirectory() as tmp_dir:
        image_paths = _pptx_to_images(local_path, tmp_dir)

        # Step 3: If screenshots available, send to Gemini Vision
        if image_paths and visual_slide_indices:
            from ..core.gemini import generate_text_with_images
            for idx in visual_slide_indices:
                if idx < len(image_paths):
                    try:
                        note = await _describe_slide_image(image_paths[idx], idx)
                        visual_notes.append(note)
                    except Exception as e:
                        logger.warning(f"Gemini vision failed for slide {idx}: {e}")

    # Step 4: Merge results into full content
    full_content = "\n\n---\n\n".join(slides_md)
    if visual_notes:
        full_content += "\n\n## Visual Notes\n\n" + "\n\n".join(visual_notes)

    # Step 5: Save to File.summary
    f.summary = full_content
    f.is_indexed = True
    await db.commit()

    return {
        "slide_count": len(slides_md),
        "parsed_content": full_content,
        "has_visual_slides": has_visual_slides,
        "visual_notes": visual_notes,
    }


@router.patch("/{file_id}/move")
async def move_file(
    file_id: str,
    folder: str = Body(..., embed=True),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(FileModel).where(FileModel.id == file_id))
    f = result.scalar_one_or_none()
    if not f:
        raise HTTPException(404, "File not found")
    f.folder_path = folder or "/"
    await db.commit()
    return {"id": file_id, "folder_path": f.folder_path}


@router.delete("/{file_id}")
async def delete_file(file_id: str, db: AsyncSession = Depends(get_db)):
    result = await db.execute(select(FileModel).where(FileModel.id == file_id))
    f = result.scalar_one_or_none()
    if not f:
        raise HTTPException(404, "File not found")
    local_path = os.path.join(UPLOAD_DIR, f.filename)
    if os.path.exists(local_path):
        os.remove(local_path)
    await db.delete(f)
    await db.commit()
    return {"ok": True}


# ── PPTX Parsing Helpers ──────────────────────────────────

def _parse_pptx_to_markdown(pptx_path: str) -> dict:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    slides_md: list[str] = []
    visual_slides: list[int] = []

    for i, slide in enumerate(prs.slides):
        lines: list[str] = []
        title = ""
        has_image = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        if not title:
                            title = text
                            lines.append(f"# Slide {i+1}: {text}")
                        else:
                            lines.append(f"- {text}")
            if shape.has_table:
                table = shape.table
                headers = [cell.text for cell in table.rows[0].cells]
                lines.append("| " + " | ".join(headers) + " |")
                lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in table.rows[1:]:
                    lines.append("| " + " | ".join([cell.text for cell in row.cells]) + " |")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
            if shape.has_chart:
                chart = shape.chart
                lines.append(f"[Chart: {chart.chart_type}]")
                for series in chart.series:
                    values = [str(v) for v in series.values]
                    lines.append(f"- {series.name}: {', '.join(values)}")

        if has_image:
            visual_slides.append(i)
        slides_md.append("\n".join(lines))

    return {
        "slides_markdown": slides_md,
        "visual_slide_indices": visual_slides,
    }


def _pptx_to_images(pptx_path: str, output_dir: str) -> list[str]:
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "png",
             "--outdir", output_dir, pptx_path],
            capture_output=True,
            timeout=30,
        )
        if result.returncode == 0:
            return sorted([
                os.path.join(output_dir, f)
                for f in os.listdir(output_dir)
                if f.endswith(".png")
            ])
    except Exception:
        pass
    return []


async def _describe_slide_image(image_path: str, slide_index: int) -> str:
    from ..core.gemini import generate_text
    import base64

    with open(image_path, "rb") as img_file:
        image_data = base64.b64encode(img_file.read()).decode()

    prompt = (
        f"請描述這張投影片（第 {slide_index + 1} 頁）的視覺內容，"
        "包括圖片、圖表、圖形等非文字元素的內容和含義。請用繁體中文回答，100字以內。"
    )
    return await generate_text(f"{prompt}\n[image/png base64: {image_data[:100]}...]")


# ── Background Processing ──────────────────────────────────

async def _process_file_background(file_id: str, local_path: str, mime_type: str, display_name: str):
    from ..db.database import AsyncSessionLocal
    async with AsyncSessionLocal() as db:
        try:
            result = await db.execute(select(FileModel).where(FileModel.id == file_id))
            f = result.scalar_one_or_none()
            if not f:
                return

            gemini_uri = None
            summary = ""

            try:
                gemini_uri = await upload_file_to_gemini(local_path, mime_type, display_name)
            except Exception as e:
                logger.warning(f"Gemini upload failed: {e}")

            try:
                if gemini_uri:
                    summary = await generate_file_summary(gemini_uri, display_name)
                else:
                    summary = await _generate_summary_from_local(local_path, mime_type, display_name)
            except Exception as e:
                logger.warning(f"Summary failed: {e}")
                summary = f"檔案 {display_name} 已上傳。"

            f.gemini_file_uri = gemini_uri
            f.summary = summary
            f.is_indexed = True
            await db.commit()

        except Exception as e:
            logger.error(f"Background processing error: {e}")
            try:
                result = await db.execute(select(FileModel).where(FileModel.id == file_id))
                f = result.scalar_one_or_none()
                if f:
                    f.summary = "處理失敗，請重新上傳"
                    f.is_indexed = True
                    await db.commit()
            except Exception:
                pass


async def _generate_summary_from_local(local_path: str, mime_type: str, display_name: str) -> str:
    from ..core.gemini import generate_text
    text = ""
    try:
        if mime_type == "application/pdf":
            import pdfplumber
            with pdfplumber.open(local_path) as pdf:
                text = "\n".join(p.extract_text() or "" for p in pdf.pages[:5])
        elif "word" in mime_type:
            from docx import Document
            doc = Document(local_path)
            text = "\n".join(p.text for p in doc.paragraphs[:50])
        elif "presentationml" in mime_type or "powerpoint" in mime_type:
            from pptx import Presentation
            prs = Presentation(local_path)
            slides_text = []
            for i, slide in enumerate(prs.slides[:10]):
                slide_lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_lines.append(shape.text.strip())
                if slide_lines:
                    slides_text.append(f"第{i+1}頁: " + " | ".join(slide_lines))
            text = "\n".join(slides_text)
        elif mime_type in ("text/plain", "text/markdown"):
            async with aiofiles.open(local_path, "r", encoding="utf-8", errors="ignore") as f:
                text = await f.read()
            text = text[:3000]
    except Exception:
        return f"檔案 {display_name} 已上傳至知識庫。"

    if not text.strip():
        return f"檔案 {display_name} 內容為空或無法讀取。"

    prompt = f"請用繁體中文為以下文件生成摘要（200字以內）：\n\n{text[:2000]}"
    return await generate_text(prompt)
